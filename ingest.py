import os
import uuid
import fitz                          # PyMuPDF — already in your requirements as pymupdf
from pptx import Presentation        # python-pptx
from docx import Document            # python-docx
import chromadb
from chromadb.utils import embedding_functions

# os.environ["HF_HUB_OFFLINE"]      = "1"
# os.environ["TRANSFORMERS_OFFLINE"] = "1"

from config import CHROMA_DB_PATH, COURSE_PDFS_PATH

# ── ChromaDB setup (same as before) ─────────────────────────────────────────
chroma_client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
embedding_fn  = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="BAAI/bge-small-en-v1.5"
)
chroma_collection = chroma_client.get_or_create_collection(
    name="course_materials",
    embedding_function=embedding_fn,
)


# ── Extractors — each returns list of {text, page_number, page_label} ───────

def extract_pdf(file_path: str) -> list:
    chunks = []
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text = page.get_text("text").strip()
            if text:
                chunks.append(text)
        doc.close()
    except Exception as e:
        print(f"  [PDF error] {file_path}: {e}")
    return chunks


def extract_pptx(file_path: str) -> list:
    chunks = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            parts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            text = "\n".join(parts).strip()
            if text:
                chunks.append(text)
    except Exception as e:
        print(f"  [PPTX error] {file_path}: {e}")
    return chunks


def extract_docx(file_path: str) -> list:
    """
    Chunks by natural paragraph boundaries.
    Starts a new chunk when it hits a heading style OR when the
    current chunk exceeds 800 chars — whichever comes first.
    This keeps related content (e.g. 'Session Layer' heading +
    its description) together in the same chunk.
    """
    chunks = []
    try:
        doc          = Document(file_path)
        current_text = ""

        for para in doc.paragraphs:
            line  = para.text.strip()
            style = (para.style.name or "").lower()

            if not line:
                continue

            is_heading = (
                "heading" in style
                or line.isupper()                    # ALL CAPS lines = likely a heading
                or (len(line) < 80 and line.endswith(":"))  # short lines ending in colon
            )

            # Start a new chunk when we hit a heading (if current chunk has content)
            # OR when current chunk is already long enough
            if current_text and (is_heading or len(current_text) >= 800):
                chunks.append(current_text.strip())
                current_text = ""

            current_text += line + "\n"

        # Don't forget the last chunk
        if current_text.strip():
            chunks.append(current_text.strip())

    except Exception as e:
        print(f"  [DOCX error] {file_path}: {e}")

    return chunks


def extract_txt(file_path: str) -> list:
    chunks = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read().strip()
        for i in range(0, len(text), 500):
            chunk = text[i:i + 500].strip()
            if chunk:
                chunks.append(chunk)
    except Exception as e:
        print(f"  [TXT error] {file_path}: {e}")
    return chunks


def extract_csv(file_path: str) -> list:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read().strip()
        return [text[:2000]] if text else []
    except Exception as e:
        print(f"  [CSV error] {file_path}: {e}")
        return []


def extract_file(file_path: str) -> list:
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf":  extract_pdf,
        ".pptx": extract_pptx,
        ".docx": extract_docx,
        ".doc":  extract_docx,
        ".txt":  extract_txt,
        ".csv":  extract_csv,
    }
    fn = extractors.get(ext)
    if fn:
        return fn(file_path)
    print(f"  [skip] Unsupported type: {file_path}")
    return []


# ── Core ingest function (called both by ingest_all and by the /upload API) ──

def ingest_file(file_path: str, course_name: str, module_name: str = "") -> int:
    file_name = os.path.basename(file_path)
    print(f"  Ingesting: {file_name}")

    chunks = extract_file(file_path)
    if not chunks:
        print(f"  [warn] No text extracted from {file_name}")
        return 0

    # Safe re-ingest — delete existing chunks for this file first
    try:
        chroma_collection.delete(where={"file_name": {"$eq": file_name}})
    except Exception:
        pass

    ids       = []
    documents = []
    metadatas = []

    for chunk in chunks:
        ids.append(str(uuid.uuid4()))
        documents.append(chunk)           # plain string now
        metadatas.append({
            "file_name": file_name,
            "file_path": file_path,
            "course":    course_name,
            "module":    module_name or "",
        })

    chroma_collection.add(ids=ids, documents=documents, metadatas=metadatas)
    print(f"  ✓ {len(chunks)} chunks indexed from {file_name}")
    return len(chunks)


# ── Walk Course_Materials and ingest everything ──────────────────────────────

def ingest_all():
    """
    Supports both flat and module-subfolder structures:

      Course_Materials/
        └── Computer Networks/
              └── CIT421.pdf              ← flat (no modules)
              └── Module 1 - Intro/
                    └── lecture.pdf       ← with module subfolders
    """
    print(f"[ingest_all] COURSE_PDFS_PATH={COURSE_PDFS_PATH}")
    if not os.path.exists(COURSE_PDFS_PATH):
        print(f"[error] Path not found: {COURSE_PDFS_PATH}")
        return 0

    # Diagnostic: show top-level listing
    try:
        top_entries = sorted(os.listdir(COURSE_PDFS_PATH))
        print(f"[ingest_all] Top-level entries ({len(top_entries)}): {top_entries}")
    except Exception as e:
        print(f"[ingest_all] Failed to list {COURSE_PDFS_PATH}: {e}")

    total = 0
    for course_folder in sorted(os.listdir(COURSE_PDFS_PATH)):
        course_path = os.path.join(COURSE_PDFS_PATH, course_folder)
        if not os.path.isdir(course_path):
            continue

        print(f"\nCourse: {course_folder}")
        # Diagnostic: list course folder contents
        try:
            items = sorted(os.listdir(course_path))
            print(f"  [debug] {course_folder} contains {len(items)} entries: {items}")
        except Exception as e:
            print(f"  [debug] Failed to list {course_path}: {e}")

        for item in sorted(os.listdir(course_path)):
            item_path = os.path.join(course_path, item)

            if os.path.isfile(item_path):
                # File sits directly in the course folder — no module
                total += ingest_file(item_path, course_folder, module_name="")

            elif os.path.isdir(item_path):
                # Module subfolder
                module_name = item
                print(f"  Module: {module_name}")
                for fname in sorted(os.listdir(item_path)):
                    fpath = os.path.join(item_path, fname)
                    if os.path.isfile(fpath):
                        # Diagnostic: file size + ext
                        try:
                            sz = os.path.getsize(fpath)
                        except Exception:
                            sz = -1
                        print(f"    [file] {fname} (size={sz} bytes)")
                        total += ingest_file(fpath, course_folder, module_name)

    if total == 0:
        print("\n[warning] Ingestion finished but indexed 0 chunks. Possible causes: empty folders, unsupported file types, or scanned/image-only PDFs without extractable text.")
    print(f"\n✓ Ingestion complete. Total chunks: {total}")
    return total


if __name__ == "__main__":
    ingest_all()


def ingest_from_url(
    file_url: str,
    file_name: str,
    course_name: str,
    module_name: str = "",
) -> int:
    """
    Downloads a file from a Cloudinary URL into a temp file,
    then calls the existing ingest_file() logic.
    Called by the POST /upload endpoint in api.py.
    """
    import tempfile
    import urllib.request

    ext = os.path.splitext(file_name)[1].lower()
    tmp_path = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp_path = tmp.name

        # Download the file from Cloudinary with proper headers
        print(f"  [ingest_from_url] Downloading {file_name} from Cloudinary...")
        print(f"  [ingest_from_url] URL: {file_url[:80]}...")
        
        import ssl
        req = urllib.request.Request(
            file_url,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        )
        
        # Handle SSL certificate issues
        ssl_context = ssl.create_default_context()
        ssl_context.check_hostname = False
        ssl_context.verify_mode = ssl.CERT_NONE
        
        with urllib.request.urlopen(req, context=ssl_context, timeout=30) as response:
            with open(tmp_path, 'wb') as f:
                f.write(response.read())
        
        print(f"  [ingest_from_url] Downloaded {file_name} successfully ({os.path.getsize(tmp_path)} bytes)")

        # Reuse existing ingest logic — pass tmp_path as file_path
        # but use the real file_name for metadata
        chunks = extract_file(tmp_path)

        if not chunks:
            print(f"  [ERROR] No text extracted from {file_name} — file may be corrupted or unsupported format")
            return 0

        # Safe re-ingest — delete existing chunks for this file first
        try:
            chroma_collection.delete(where={"file_name": {"$eq": file_name}})
        except Exception:
            pass

        ids, documents, metadatas = [], [], []
        for chunk in chunks:
            ids.append(str(uuid.uuid4()))
            documents.append(chunk)
            metadatas.append({
                "file_name":   file_name,
                "file_url":    file_url,       # store Cloudinary URL in metadata
                "course":      course_name,
                "module":      module_name or "",
            })

        chroma_collection.add(ids=ids, documents=documents, metadatas=metadatas)
        print(f"  ✓ {len(chunks)} chunks indexed from {file_name}")
        return len(chunks)

    except urllib.error.HTTPError as e:
        print(f"  [ERROR] HTTP Error {e.code} downloading {file_name}: {e.reason}")
        if e.code == 401:
            print(f"  [ERROR] URL unauthorized - Cloudinary URL may be expired or private")
        elif e.code == 404:
            print(f"  [ERROR] File not found at URL")
        print(f"  [ERROR] Full URL: {file_url}")
        return 0
    except urllib.error.URLError as e:
        print(f"  [ERROR] Failed to download {file_name}: {e.reason}")
        return 0
    except TimeoutError as e:
        print(f"  [ERROR] Download timeout for {file_name}: {e}")
        return 0
    except Exception as e:
        print(f"  [ERROR] ingest_from_url failed for {file_name}: {type(e).__name__}: {e}")
        import traceback
        traceback.print_exc()
        return 0

    finally:
        # Always clean up the temp file
        if tmp_path:
            try:
                os.unlink(tmp_path)
                print(f"  [ingest_from_url] Cleaned up temp file: {tmp_path}")
            except Exception as e:
                print(f"  [warn] Failed to clean temp file {tmp_path}: {e}")


def delete_file(file_name: str) -> dict:
    """
    Removes all ChromaDB chunks associated with a specific file.
    Called by the DELETE /material/{file_name} endpoint in api.py.
    """
    try:
        # Count how many chunks exist for this file before deleting
        existing = chroma_collection.get(
            where={"file_name": {"$eq": file_name}},
            include=["metadatas"]
        )
        chunk_count = len(existing.get("metadatas", []))

        if chunk_count == 0:
            return {
                "success": False,
                "message": f"No chunks found for '{file_name}'. It may not have been ingested."
            }

        chroma_collection.delete(
            where={"file_name": {"$eq": file_name}}
        )

        return {
            "success": True,
            "message": f"Deleted {chunk_count} chunks for '{file_name}' from ChromaDB."
        }

    except Exception as e:
        return {
            "success": False,
            "message": f"Delete failed: {str(e)}"
        }